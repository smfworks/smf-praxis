"""Deterministic core and optional professional renderer contracts."""
from __future__ import annotations

import base64
import builtins
import io
import zipfile
from dataclasses import replace
from pathlib import Path
from types import MappingProxyType

import pytest

from hybridagent.artifacts import (
    ArtifactRenderError,
    MissingArtifactBackendError,
    ParagraphBlock,
    render_artifact,
    supported_formats,
)
from hybridagent.artifacts.render_common import (
    checked_assets,
    image_kind,
    normalize_zip_package,
)
from hybridagent.artifacts.render_docx import render_docx
from hybridagent.artifacts.render_json import render_json
from hybridagent.artifacts.render_markdown import render_markdown
from hybridagent.artifacts.render_pdf import render_pdf
from hybridagent.artifacts.render_pptx import render_pptx
from hybridagent.artifacts.render_xlsx import render_xlsx
from tests.artifact_helpers import PNG
from tests.test_artifact_models import document

GOLDEN = Path(__file__).parent / "golden"
OPTIONAL_MODULES = {"docx", "reportlab", "pptx", "openpyxl", "PIL", "pypdf"}
JPEG = base64.b64decode(
    "/9j/4AAQSkZJRgABAQAAAQABAAD/2wBDAAgGBgcGBQgHBwcJCQgKDBQNDAsLDBkSEw8UHRofHh0a"
    "HBwgJC4nICIsIxwcKDcpLDAxNDQ0Hyc5PTgyPC4zNDL/2wBDAQkJCQwLDBgNDRgyIRwhMjIyMjIy"
    "MjIyMjIyMjIyMjIyMjIyMjIyMjIyMjIyMjIyMjIyMjIyMjIyMjIyMjIyMjL/wAARCAABAAEDASIA"
    "AhEBAxEB/8QAHwAAAQUBAQEBAQEAAAAAAAAAAAECAwQFBgcICQoL/8QAtRAAAgEDAwIEAwUFBAQAAA"
    "F9AQIDAAQRBRIhMUEGE1FhByJxFDKBkaEII0KxwRVS0fAkM2JyggkKFhcYGRolJicoKSo0NTY3ODk6"
    "Q0RFRkdISUpTVFVWV1hZWmNkZWZnaGlqc3R1dnd4eXqDhIWGh4iJipKTlJWWl5iZmqKjpKWmp6ipqr"
    "KztLW2t7i5usLDxMXGx8jJytLT1NXW19jZ2uHi4+Tl5ufo6erx8vP09fb3+Pn6/8QAHwEAAwEBAQEB"
    "AQEBAQAAAAAAAAECAwQFBgcICQoL/8QAtREAAgECBAQDBAcFBAQAAQJ3AAECAxEEBSExBhJBUQdhcR"
    "MiMoEIFEKRobHBCSMzUvAVYnLRChYkNOEl8RcYGRomJygpKjU2Nzg5OkNERUZHSElKU1RVVldYWVpj"
    "ZGVmZ2hpanN0dXZ3eHl6goOEhYaHiImKkpOUlZaXmJmaoqOkpaanqKmqsrO0tba3uLm6wsPExcbHyM"
    "nK0tPU1dbX2Nna4uPk5ebn6Onq8vP09fb3+Pn6/9oADAMBAAIRAxEAPwD3+iiigD//2Q=="
)
VALID_SVG = (
    b'<?xml version="1.0"?><svg xmlns="http://www.w3.org/2000/svg" width="1" '
    b'height="1"><rect width="1" height="1"/></svg>'
)


def _png() -> bytes:
    image_module = pytest.importorskip("PIL.Image")
    output = io.BytesIO()
    image_module.new("RGB", (32, 24), (24, 72, 120)).save(output, format="PNG")
    return output.getvalue()


def _multi_page_document():
    base = document()
    extra = tuple(ParagraphBlock(f"long-{index}", f"Governed detail paragraph {index}. " * 8)
                  for index in range(80))
    first = base.sections[0]
    return replace(base, sections=(replace(first, blocks=(*first.blocks, *extra)),))


def test_core_import_does_not_load_optional_backends(monkeypatch: pytest.MonkeyPatch) -> None:
    real_import = builtins.__import__

    def guarded(name: str, *args: object, **kwargs: object) -> object:
        if name.split(".")[0] in OPTIONAL_MODULES:
            raise AssertionError(f"optional backend imported by core: {name}")
        return real_import(name, *args, **kwargs)

    monkeypatch.setattr(builtins, "__import__", guarded)
    import hybridagent.artifacts as artifacts

    assert artifacts.supported_formats() == ("json", "markdown", "docx", "pdf", "pptx", "xlsx")


def test_json_and_markdown_are_byte_deterministic_and_match_goldens() -> None:
    artifact = document()
    assert render_json(artifact) == render_json(artifact) == artifact.canonical_bytes()
    assert render_markdown(artifact) == render_markdown(artifact)
    assert render_json(artifact) == (GOLDEN / "artifact-sample.json").read_bytes()
    assert render_markdown(artifact) == (GOLDEN / "artifact-sample.md").read_bytes()


def test_markdown_contains_exact_citations_sources_and_classification() -> None:
    text = render_markdown(document()).decode()
    assert "confidentiality: \"confidential\"" in text
    assert "The supported finding. [^citation-1]" in text
    assert "source `source-1`; spans: span-1; claims: claim-1; p. 2" in text
    assert "# Source Manifest" in text
    assert "# Signatures" in text


def test_registry_is_strict_and_core_renderers_need_no_asset_bytes() -> None:
    artifact = document()
    assert supported_formats() == ("json", "markdown", "docx", "pdf", "pptx", "xlsx")
    assert render_artifact(artifact, "json") == artifact.canonical_bytes()
    assert render_artifact(artifact, "markdown") == render_markdown(artifact)
    with pytest.raises(ArtifactRenderError, match="unsupported"):
        render_artifact(artifact, "html")
    with pytest.raises(ArtifactRenderError, match="exact dictionary"):
        render_artifact(artifact, "json", MappingProxyType({}))


def test_missing_optional_backend_has_actionable_error(monkeypatch: pytest.MonkeyPatch) -> None:
    real_import = builtins.__import__

    def missing(name: str, *args: object, **kwargs: object) -> object:
        if name.split(".")[0] == "docx":
            raise ImportError("absent")
        return real_import(name, *args, **kwargs)

    monkeypatch.setattr(builtins, "__import__", missing)
    with pytest.raises(
        MissingArtifactBackendError,
        match=r'from a Praxis source checkout.*pip install -e "\.\[artifacts\]"',
    ):
        render_docx(document(), {"asset-1": b"not-used"})


def test_office_zip_normalization_rejects_nonportable_and_case_colliding_paths() -> None:
    for members in (
        (("C:/absolute.xml", b"x"),),
        (("word/CON.xml", b"x"),),
        (("word/A.xml", b"x"), ("word/a.xml", b"y")),
    ):
        source = io.BytesIO()
        with zipfile.ZipFile(source, "w") as archive:
            for name, payload in members:
                archive.writestr(name, payload)
        with pytest.raises(ArtifactRenderError, match="portable|collide"):
            normalize_zip_package(source.getvalue())


def test_optional_renderers_reject_missing_or_malformed_assets() -> None:
    pytest.importorskip("docx")
    with pytest.raises(ArtifactRenderError, match="missing"):
        render_docx(document(), {})
    with pytest.raises(ArtifactRenderError, match="recognized|declared media type"):
        render_docx(document(), {"asset-1": b"not-an-image"})


@pytest.mark.parametrize(
    ("media_type", "payload"),
    (
        ("image/png", b"\x89PNG\r\n\x1a\n"),
        ("image/jpeg", b"\xff\xd8\xff"),
        (
            "image/jpeg",
            b"\xff\xd8\xff\xc0\x00\x0b\x08\x00\x01\x00\x01\x01\x01\x11\x00"
            b"\xff\xda\x00\x08\x01\x01\x00\x00\x3f\x00\xff\xd9",
        ),
        (
            "image/jpeg",
            b"\xff\xd8\xff\xda\x00\x08\x01\x01\x00\x00\x3f\x00\x00"
            b"\xff\xc0\x00\x0b\x08\x00\x01\x00\x01\x01\x01\x11\x00\xff\xd9",
        ),
        ("image/svg+xml", b"<svg"),
        ("image/svg+xml", b"prefix <svg xmlns='http://www.w3.org/2000/svg'/>")
    ),
)
def test_asset_admission_rejects_signature_only_or_malformed_images(
    media_type: str, payload: bytes
) -> None:
    artifact = document()
    section = artifact.sections[0]
    figure = replace(section.blocks[3], media_type=media_type)
    artifact = replace(
        artifact,
        sections=(replace(section, blocks=(*section.blocks[:3], figure)),),
    )
    with pytest.raises(ArtifactRenderError, match="declared media type"):
        checked_assets(artifact, {"asset-1": payload})


def test_asset_admission_accepts_complete_png_jpeg_and_svg() -> None:
    assert image_kind(PNG) == "png"
    assert image_kind(JPEG) == "jpeg"
    assert image_kind(VALID_SVG) == "svg"


def test_all_optional_renderers_are_deterministic_and_semantically_valid() -> None:
    pytest.importorskip("docx")
    pytest.importorskip("reportlab")
    pytest.importorskip("pptx")
    pytest.importorskip("openpyxl")
    pytest.importorskip("pypdf")
    from docx import Document
    from openpyxl import load_workbook
    from pptx import Presentation
    from pypdf import PdfReader

    artifact = document()
    assets = {"asset-1": _png()}
    rendered = {
        "docx": render_docx(artifact, assets),
        "pdf": render_pdf(artifact, assets),
        "pptx": render_pptx(artifact, assets),
        "xlsx": render_xlsx(artifact, assets),
    }
    assert rendered == {
        "docx": render_docx(artifact, assets),
        "pdf": render_pdf(artifact, assets),
        "pptx": render_pptx(artifact, assets),
        "xlsx": render_xlsx(artifact, assets),
    }

    docx = Document(io.BytesIO(rendered["docx"]))
    assert docx.core_properties.title == artifact.metadata.title
    assert artifact.metadata.title in "\n".join(p.text for p in docx.paragraphs)
    with zipfile.ZipFile(io.BytesIO(rendered["docx"])) as archive:
        docx_xml = b"\n".join(archive.read(name) for name in archive.namelist() if name.endswith(".xml"))
    assert b' descr="Inspection overview"' in docx_xml
    assert b"w:tblHeader" in docx_xml

    pdf = PdfReader(io.BytesIO(rendered["pdf"]))
    assert pdf.metadata.title == artifact.metadata.title
    pdf_text = "\n".join(page.extract_text() or "" for page in pdf.pages)
    assert artifact.metadata.title in pdf_text
    assert "Source Manifest" in pdf_text

    pptx = Presentation(io.BytesIO(rendered["pptx"]))
    pptx_text = "\n".join(
        shape.text for slide in pptx.slides for shape in slide.shapes if hasattr(shape, "text")
    )
    assert artifact.metadata.title in pptx_text
    with zipfile.ZipFile(io.BytesIO(rendered["pptx"])) as archive:
        pptx_xml = b"\n".join(archive.read(name) for name in archive.namelist() if name.endswith(".xml"))
    assert b' descr="Inspection overview"' in pptx_xml

    xlsx = load_workbook(io.BytesIO(rendered["xlsx"]), data_only=False)
    assert xlsx.properties.title == artifact.metadata.title
    assert {"Overview", "Sections", "Sources", "Citations", "Governance"} <= set(xlsx.sheetnames)
    assert xlsx["Overview"]["B2"].value == artifact.artifact_id

    for format_name in ("docx", "pptx", "xlsx"):
        with zipfile.ZipFile(io.BytesIO(rendered[format_name])) as archive:
            assert len(archive.namelist()) == len(set(archive.namelist()))
            assert all(item.date_time == (1980, 1, 1, 0, 0, 0) for item in archive.infolist())
            assert all(not name.startswith("/") and ".." not in name.split("/")
                       for name in archive.namelist())


def test_pdf_pagination_is_stable_and_footer_is_present() -> None:
    pytest.importorskip("reportlab")
    pytest.importorskip("pypdf")
    from pypdf import PdfReader

    artifact = _multi_page_document()
    payload = render_pdf(artifact, {"asset-1": _png()})
    reader = PdfReader(io.BytesIO(payload))
    assert len(reader.pages) >= 2
    text = "\n".join(page.extract_text() or "" for page in reader.pages)
    assert artifact.artifact_id in text
    assert artifact.metadata.confidentiality.upper() in text


def test_renderer_rejects_stateful_mapping_subclasses() -> None:
    class Assets(dict[str, bytes]):
        pass

    with pytest.raises(ArtifactRenderError, match="exact dictionary"):
        render_artifact(document(), "docx", Assets({"asset-1": b"x"}))


def test_renderers_neutralize_markup_and_spreadsheet_formula_injection() -> None:
    pytest.importorskip("reportlab")
    pytest.importorskip("pypdf")
    pytest.importorskip("openpyxl")
    from openpyxl import load_workbook
    from pypdf import PdfReader

    base = document()
    section = base.sections[0]
    hostile = replace(
        base,
        metadata=replace(base.metadata, title='Opinion "quoted" <script>alert(1)</script>'),
        sections=(replace(section, blocks=(
            ParagraphBlock("hostile-paragraph", "<b>literal evidence</b> & not markup"),
            replace(section.blocks[2], rows=(("A", '=HYPERLINK("https://evil.test")'),)),
            section.blocks[3],
        )),),
    )
    markdown = render_markdown(hostile).decode()
    assert "&lt;b&gt;literal evidence&lt;/b&gt; &amp; not markup" in markdown
    assert "# Opinion \"quoted\" &lt;script&gt;alert(1)&lt;/script&gt;" in markdown

    assets = {"asset-1": _png()}
    pdf = PdfReader(io.BytesIO(render_pdf(hostile, assets)))
    pdf_text = "\n".join(page.extract_text() or "" for page in pdf.pages)
    assert "<b>literal evidence</b> & not markup" in pdf_text

    workbook = load_workbook(io.BytesIO(render_xlsx(hostile, assets)), data_only=False)
    formula_cells = [cell for sheet in workbook.worksheets for row in sheet.iter_rows()
                     for cell in row if cell.data_type == "f"]
    assert formula_cells == []
    assert workbook["Results"]["B2"].value.startswith("'=HYPERLINK")
