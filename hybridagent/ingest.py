"""Document ingestion — extract text from PDF, Office, email, HTML, and plain
files so they can be chunked, embedded, and retrieved (RAG).

Design goals:

* **Core stays dependency-free.** Plain text, Markdown, CSV/JSON/logs, ``.eml``
  email, and basic HTML parse with the standard library alone.
* **Heavier formats are optional.** PDF/DOCX/PPTX/XLSX/MSG import their parser
  lazily; if the dependency is missing we raise :class:`MissingDependencyError`
  naming the exact extra to install (``pip install "praxis-agent[docs]"``).
* **Best parser first.** If Microsoft's ``markitdown`` is installed it is used
  for the rich formats (it yields clean Markdown); otherwise per-format
  fallbacks run.

Everything returns an :class:`ExtractedDoc` with provenance so the governance and
memory layers can keep treating retrieved content as *data, never instruction*.
"""
from __future__ import annotations

import csv
import io
import json
from dataclasses import dataclass, field
from html.parser import HTMLParser
from pathlib import Path

from .evidence import EvidenceRegistry, EvidenceSource, EvidenceVersion
from .extraction import EvidenceSpan, ExtractionRegistry
from .logging_util import get_logger
from .persistence import Store

_log = get_logger("praxis.ingest")


class MissingDependencyError(RuntimeError):
    """Raised when an optional parser dependency is not installed."""


@dataclass
class ExtractedDoc:
    text: str
    source: str
    kind: str = "document"
    metadata: dict = field(default_factory=dict)


def register_evidence(
    document: ExtractedDoc, original: bytes, *, store: Store,
    organization_id: str, workspace_id: str, created_by: str,
    canonical_uri: str, publisher: str, locator: dict,
) -> tuple[EvidenceSource, EvidenceVersion, EvidenceSpan]:
    """Persist one extracted document through canonical evidence lineage."""
    evidence = EvidenceRegistry(store)
    source = evidence.create_source(
        organization_id, workspace_id, canonical_uri=canonical_uri,
        publisher=publisher, created_by=created_by)
    suffix = Path(document.source).suffix.lower()
    mime_type = {
        ".txt": "text/plain", ".md": "text/markdown", ".pdf": "application/pdf",
        ".html": "text/html", ".json": "application/json",
    }.get(suffix, "application/octet-stream")
    version = evidence.add_version(
        organization_id, workspace_id, source.source_id, content=original,
        mime_type=mime_type, retrieved_ts=float(document.metadata.get("retrieved_ts", 0.0)),
        parser=str(document.metadata.get("parser", document.kind)),
        parser_version=str(document.metadata.get("parser_version", "1")),
        parser_config=dict(document.metadata.get("parser_config", {})),
        license=str(document.metadata.get("license", "unspecified")),
        original_object_path=str(document.metadata.get("path", document.source)),
        created_by=created_by)
    span = ExtractionRegistry(store).add_span(
        organization_id, workspace_id, version.version_id,
        locator_type="document", locator=locator,
        extracted_text=document.text, created_by=created_by)
    return source, version, span


TEXT_SUFFIXES = {".txt", ".md", ".markdown", ".log", ".rst", ".csv", ".tsv", ".json"}
RICH_SUFFIXES = {".pdf", ".docx", ".pptx", ".xlsx", ".msg", ".html", ".htm", ".eml"}
SUPPORTED = TEXT_SUFFIXES | RICH_SUFFIXES


def is_supported(path: str | Path) -> bool:
    return Path(path).suffix.lower() in SUPPORTED


# --------------------------------------------------------------------- helpers
class _TextHTMLParser(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self._parts: list[str] = []
        self._skip = False

    def handle_starttag(self, tag, attrs):
        if tag in ("script", "style"):
            self._skip = True

    def handle_endtag(self, tag):
        if tag in ("script", "style"):
            self._skip = False

    def handle_data(self, data):
        if not self._skip and data.strip():
            self._parts.append(data.strip())

    def text(self) -> str:
        return "\n".join(self._parts)


def _missing(dep: str, extra: str = "docs"):
    raise MissingDependencyError(
        f"'{dep}' is required for this file type. Install with: "
        f'pip install "praxis-agent[{extra}]"  (or pip install {dep})')


def _try_markitdown(path: Path) -> str | None:
    try:
        from markitdown import MarkItDown  # type: ignore
    except Exception:
        return None
    try:
        md = MarkItDown()
        return md.convert(str(path)).text_content
    except Exception as exc:  # fall back to per-format parser
        _log.warning("markitdown failed on %s (%s); using fallback", path.name, exc)
        return None


# --------------------------------------------------------------------- parsers
def _read_text(path: Path) -> str:
    suffix = path.suffix.lower()
    raw = path.read_text(encoding="utf-8", errors="replace")
    if suffix in (".csv", ".tsv"):
        delim = "\t" if suffix == ".tsv" else ","
        # csv.reader raises "_csv.Error: line contains NUL" on embedded NUL bytes
        # on Python < 3.11 (3.11+ tolerates them). Strip NULs first so a poisoned
        # file degrades gracefully and parses identically on every supported
        # Python version instead of crashing the whole ingestion run.
        rows = list(csv.reader(io.StringIO(raw.replace("\x00", "")), delimiter=delim))
        return "\n".join(", ".join(r) for r in rows)
    if suffix == ".json":
        try:
            return json.dumps(json.loads(raw), indent=2)
        except json.JSONDecodeError:
            pass
    return raw


def _parse_eml(path: Path) -> str:
    from email import policy
    from email.parser import BytesParser
    with path.open("rb") as fh:
        msg = BytesParser(policy=policy.default).parse(fh)
    head = (f"From: {msg.get('from','')}\nTo: {msg.get('to','')}\n"
            f"Date: {msg.get('date','')}\nSubject: {msg.get('subject','')}\n\n")
    body = ""
    try:
        part = msg.get_body(preferencelist=("plain", "html"))
        if part is not None:
            content = str(part.get_content())
            if part.get_content_subtype() == "html":
                p = _TextHTMLParser()
                p.feed(content)
                content = p.text()
            body = content
    except Exception:
        payload = msg.get_payload(decode=False)
        body = payload if isinstance(payload, str) else ""
    return head + body


def _parse_html(path: Path) -> str:
    p = _TextHTMLParser()
    p.feed(path.read_text(encoding="utf-8", errors="replace"))
    return p.text()


def _parse_pdf(path: Path) -> str:
    try:
        from pypdf import PdfReader  # type: ignore
    except Exception:
        _missing("pypdf")
    reader = PdfReader(str(path))
    return "\n\n".join((page.extract_text() or "") for page in reader.pages)


def _parse_docx(path: Path) -> str:
    try:
        import docx  # type: ignore[import-untyped]
    except Exception:
        _missing("python-docx")
    d = docx.Document(str(path))
    parts = [p.text for p in d.paragraphs if p.text.strip()]
    for table in d.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _parse_pptx(path: Path) -> str:
    try:
        from pptx import Presentation  # type: ignore
    except Exception:
        _missing("python-pptx")
    prs = Presentation(str(path))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"# Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text)
    return "\n".join(parts)


def _parse_xlsx(path: Path) -> str:
    try:
        import openpyxl  # type: ignore
    except Exception:
        _missing("openpyxl")
    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    parts: list[str] = []
    for ws in wb.worksheets:
        parts.append(f"# Sheet: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                parts.append(", ".join(cells))
    return "\n".join(parts)


def _parse_msg(path: Path) -> str:
    try:
        import extract_msg  # type: ignore
    except Exception:
        _missing("extract-msg")
    m = extract_msg.Message(str(path))
    try:
        return (f"From: {m.sender or ''}\nTo: {m.to or ''}\nDate: {m.date or ''}\n"
                f"Subject: {m.subject or ''}\n\n{m.body or ''}")
    finally:
        try:
            m.close()
        except Exception:
            pass


_PARSERS = {
    ".eml": _parse_eml, ".html": _parse_html, ".htm": _parse_html,
    ".pdf": _parse_pdf, ".docx": _parse_docx, ".pptx": _parse_pptx,
    ".xlsx": _parse_xlsx, ".msg": _parse_msg,
}


# ------------------------------------------------------------------- front door
def extract_text(path: str | Path) -> ExtractedDoc:
    p = Path(path)
    if not p.exists():
        raise FileNotFoundError(p)
    suffix = p.suffix.lower()
    kind = {
        ".pdf": "pdf", ".docx": "word", ".pptx": "powerpoint", ".xlsx": "excel",
        ".eml": "email", ".msg": "email", ".html": "web", ".htm": "web",
    }.get(suffix, "document")

    if suffix in TEXT_SUFFIXES:
        text = _read_text(p)
    elif suffix in _PARSERS:
        # markitdown first for the rich binary/office/web formats (not eml).
        text = None
        if suffix not in (".eml",):
            text = _try_markitdown(p)
        if text is None:
            text = _PARSERS[suffix](p)
    else:
        raise ValueError(f"unsupported file type: {suffix or '(none)'}")

    return ExtractedDoc(text=text or "", source=p.name, kind=kind,
                        metadata={"path": str(p), "suffix": suffix})
