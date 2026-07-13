"""Optional PPTX renderer for professional artifact briefings."""
from __future__ import annotations

import io
from datetime import datetime
from typing import Any

from hybridagent.artifacts.models import (
    ArtifactDocument,
    FigureBlock,
    ListBlock,
    ParagraphBlock,
    TableBlock,
)
from hybridagent.artifacts.render_common import (
    ArtifactRenderError,
    checked_assets,
    image_kind,
    normalize_zip_package,
    require_backend,
)

_FIXED_TIME = datetime(2000, 1, 1)


def render_pptx(
    document: ArtifactDocument, assets: dict[str, bytes] | None = None
) -> bytes:
    require_backend("pptx", "python-pptx")
    from pptx import Presentation  # type: ignore[import-untyped]
    from pptx.util import Inches, Pt  # type: ignore[import-untyped]

    resolved = checked_assets(document, assets)
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    props = presentation.core_properties
    props.title = document.metadata.title
    props.subject = document.metadata.subject
    props.author = document.metadata.created_by
    props.last_modified_by = "Praxis Artifact Studio"
    props.created = _FIXED_TIME
    props.modified = _FIXED_TIME
    props.comments = f"artifact={document.artifact_id}; sha256={document.content_hash()}"

    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = document.metadata.title
    subtitle = title_slide.placeholders[1]
    subtitle.text = (
        f"{document.metadata.document_type}\n"
        f"{document.metadata.confidentiality.upper()}\n"
        f"Artifact {document.artifact_id}"
    )

    def content_slide(title: str) -> Any:
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = title
        return slide

    for section in (*document.sections, *document.appendices):
        slide = content_slide(section.title)
        body = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(11.9), Inches(5.3))
        frame = body.text_frame
        frame.clear()
        frame.word_wrap = True
        first = True
        for block in section.blocks:
            if type(block) is ParagraphBlock:
                paragraph = frame.paragraphs[0] if first else frame.add_paragraph()
                first = False
                paragraph.text = block.text
                paragraph.level = 0
                paragraph.font.size = Pt(18)
            elif type(block) is ListBlock:
                for item in block.items:
                    paragraph = frame.paragraphs[0] if first else frame.add_paragraph()
                    first = False
                    paragraph.text = item
                    paragraph.level = 0
                    paragraph.font.size = Pt(18)
            elif type(block) is TableBlock:
                table_slide = content_slide(block.caption or section.title)
                shape = table_slide.shapes.add_table(
                    len(block.rows) + 1, len(block.columns),
                    Inches(0.6), Inches(1.4), Inches(12.1), Inches(5.3),
                )
                table = shape.table
                for col, value in enumerate(block.columns):
                    table.cell(0, col).text = value
                for row_index, row in enumerate(block.rows, 1):
                    for col, value in enumerate(row):
                        table.cell(row_index, col).text = value
            elif type(block) is FigureBlock:
                payload = resolved[block.asset_id]
                if image_kind(payload) == "svg":
                    raise ArtifactRenderError("PPTX rendering requires PNG or JPEG figure assets")
                figure_slide = content_slide(block.caption or section.title)
                picture = figure_slide.shapes.add_picture(
                    io.BytesIO(payload), Inches(1.0), Inches(1.5),
                    width=Inches(8.5), height=Inches(4.8),
                )
                picture._pic.nvPicPr.cNvPr.set("descr", block.alt_text)
                picture._pic.nvPicPr.cNvPr.set("title", block.caption or block.alt_text)
                caption = figure_slide.shapes.add_textbox(
                    Inches(9.8), Inches(2.0), Inches(2.8), Inches(3.0)
                )
                caption.text_frame.text = block.caption or block.alt_text
        if first:
            frame.paragraphs[0].text = "Section content is presented in the following slides."

    source_slide = content_slide("Source Manifest")
    source_box = source_slide.shapes.add_textbox(
        Inches(0.7), Inches(1.4), Inches(11.9), Inches(5.3)
    )
    source_frame = source_box.text_frame
    source_frame.clear()
    for index, source in enumerate(document.sources):
        paragraph = source_frame.paragraphs[0] if index == 0 else source_frame.add_paragraph()
        paragraph.text = f"{source.source_id} — {source.title or source.source_version_id} — {source.content_hash}"
        paragraph.font.size = Pt(12)
    if not document.sources:
        source_frame.paragraphs[0].text = "No source manifest entries."

    output = io.BytesIO()
    presentation.save(output)
    return normalize_zip_package(output.getvalue())
